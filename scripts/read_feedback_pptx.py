
from pptx import Presentation
import os

FILE_PATH = r"c:\Users\jlcue\Documents\mag-sistema\ref\COMENTARIOS A PANTALLAS 26 DE MARZO.pptx"

if not os.path.exists(FILE_PATH):
    print(f"Error: {FILE_PATH} not found.")
else:
    prs = Presentation(FILE_PATH)
    print(f"Slides count: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- SLIDE {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"TEXT: {shape.text.strip()}")
        if slide.has_notes_slide:
            print(f"NOTES: {slide.notes_slide.notes_text_frame.text.strip()}")
